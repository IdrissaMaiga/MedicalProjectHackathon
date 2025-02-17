import fitz  # PyMuPDF for PDFs
import pytesseract
from PIL import Image
from io import BytesIO
from docx import Document  # For .docx files
from pptx import Presentation  # For .pptx files
import requests
from bs4 import BeautifulSoup
import re
from youtube_transcript_api import YouTubeTranscriptApi
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"  # Adjust for your OS

# Extract text from PDFs
def extract_text_from_pdf_bytes(pdf_bytes: bytes) -> str:
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    return "\n".join(page.get_text() for page in doc)

# Extract text from images using Tesseract OCR
def extract_text_from_image_bytes(image_bytes: bytes) -> str:
    img = Image.open(BytesIO(image_bytes))
    return pytesseract.image_to_string(img)

# Extract text from Word documents (.docx)
def extract_text_from_docx_bytes(docx_bytes: bytes) -> str:
    doc = Document(BytesIO(docx_bytes))
    return "\n".join(para.text for para in doc.paragraphs)

# Extract text from PowerPoint files (.pptx)
def extract_text_from_pptx_bytes(pptx_bytes: bytes) -> str:
    presentation = Presentation(BytesIO(pptx_bytes))
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# Extract text from plain text (.txt)
def extract_text_from_txt_bytes(txt_bytes: bytes) -> str:
    return txt_bytes.decode("utf-8")

# Extract text from a webpage (web scraping)
def extract_text_from_url(url: str) -> str:
    response = requests.get(url)
    if response.status_code == 200:
        soup = BeautifulSoup(response.content, "html.parser")
        return soup.get_text()
    return "Failed to extract text"

# Extract YouTube transcript
def extract_youtube_transcript(video_url: str) -> str:
    # Regular expression to extract video ID from various types of YouTube links
    video_id_match = re.search(r"v=([a-zA-Z0-9_-]+)", video_url)
    
    if video_id_match:
        video_id = video_id_match.group(1)
    else:
        return "Invalid YouTube URL or video ID not found"
    
    try:
        # Attempt to fetch the transcript in English first, then fallback to French
        transcript = YouTubeTranscriptApi.get_transcript(video_id, languages=['en', 'fr'])
        return "\n".join([entry["text"] for entry in transcript])
    except Exception as e:
        return f"Transcript not available. Error: {str(e)}"